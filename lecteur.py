import streamlit as st
import json
import re
import os
import hashlib
import pickle
from collections import Counter
from datetime import datetime
from io import BytesIO

import numpy as np
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from mistralai import Mistral
from fpdf import FPDF

# ============================================================
# CONFIGURATION
# ============================================================

st.set_page_config(page_title="Insight PDF Pro", page_icon="✨", layout="wide")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Google+Sans:wght@400;500;700&display=swap');
html, body, [class*="css"] { font-family: 'Google Sans', sans-serif; }
.stButton>button {
    border-radius: 20px;
    background-color: #f0f4f9;
    border: none;
    color: #1a73e8;
    font-weight: 500;
}
[data-testid="stSidebar"] {
    background-color: #f8f9fa;
    border-right: 1px solid #e1e3e1;
}
[data-testid="stChatMessage"] {
    border-radius: 16px;
    margin-bottom: 8px;
    padding: 4px 8px;
}
</style>
""", unsafe_allow_html=True)

# ============================================================
# PROMPT ANTI-HALLUCINATION
# ============================================================

SYSTEM_PROMPT = (
    "Tu es un assistant expert en analyse de documents. "
    "Réponds UNIQUEMENT en utilisant les informations du CONTEXTE fourni. "
    "Quand tu utilises une information, cite le passage source entre guillemets. "
    "Si la réponse ne se trouve pas dans le contexte, réponds exactement : "
    "'Je suis désolé, mais cette information n'est pas présente dans le document fourni.' "
    "Ne réponds jamais en utilisant tes connaissances générales si le sujet est absent du document."
)

# Dossier de cache persistant pour les embeddings
CACHE_DIR = ".embedding_cache"
os.makedirs(CACHE_DIR, exist_ok=True)

# ============================================================
# AMÉLIORATION 1 — PARSING : pdfplumber (remplace PyPDF2)
# Meilleur sur PDF complexes (tableaux, colonnes, mise en page)
# ============================================================

def extract_pdf_data(pdf_file) -> dict:
    """
    Extraction robuste avec pdfplumber.
    Fallback automatique vers PyPDF2 si pdfplumber échoue.
    """
    pdf_bytes = pdf_file.read()

    # Tentative pdfplumber (meilleure qualité)
    try:
        import pdfplumber
        pages_text = {}
        with pdfplumber.open(BytesIO(pdf_bytes)) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    pages_text[i + 1] = text
        if pages_text:
            return pages_text
    except ImportError:
        st.warning("⚠️ pdfplumber non installé. Fallback PyPDF2. `pip install pdfplumber` recommandé.")
    except Exception as e:
        st.warning(f"⚠️ pdfplumber a échoué ({e}), fallback PyPDF2.")

    # Fallback PyPDF2
    import PyPDF2
    reader = PyPDF2.PdfReader(BytesIO(pdf_bytes))
    pages_text = {}
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages_text[i + 1] = text
    return pages_text


# ============================================================
# AMÉLIORATION 2 — CACHE EMBEDDINGS PERSISTANTS (disque)
# Évite le recalcul à chaque upload du même fichier
# ============================================================

def get_cache_path(file_key: str) -> str:
    h = hashlib.md5(file_key.encode()).hexdigest()
    return os.path.join(CACHE_DIR, f"{h}.pkl")


def load_cached_embeddings(file_key: str):
    path = get_cache_path(file_key)
    if os.path.exists(path):
        try:
            with open(path, "rb") as f:
                return pickle.load(f)
        except Exception:
            return None
    return None


def save_cached_embeddings(file_key: str, chunks: list):
    path = get_cache_path(file_key)
    try:
        with open(path, "wb") as f:
            pickle.dump(chunks, f)
    except Exception as e:
        st.warning(f"⚠️ Impossible de sauvegarder le cache : {e}")


# ============================================================
# MODÈLES (mis en cache Streamlit)
# ============================================================

@st.cache_resource(show_spinner="Chargement du modèle d'embeddings…")
def load_embedding_model():
    try:
        from sentence_transformers import SentenceTransformer
        return SentenceTransformer("all-MiniLM-L6-v2")
    except ImportError:
        st.warning("⚠️ sentence-transformers non installé. Fallback sur BM25.")
        return None


@st.cache_resource(show_spinner="Chargement du reranker…")
def load_reranker():
    try:
        from sentence_transformers import CrossEncoder
        return CrossEncoder("cross-encoder/ms-marco-MiniLM-L-6-v2")
    except ImportError:
        return None


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    norm_a = np.linalg.norm(a)
    norm_b = np.linalg.norm(b)
    if norm_a == 0 or norm_b == 0:
        return 0.0
    return float(np.dot(a, b) / (norm_a * norm_b))


def encode_chunks(chunks: list, model, file_key: str) -> list:
    """
    Encode les chunks avec cache persistant.
    Si déjà calculé pour ce fichier → charge depuis le disque.
    """
    cached = load_cached_embeddings(file_key)
    if cached is not None:
        st.info("⚡ Embeddings chargés depuis le cache (aucun recalcul).")
        return cached

    texts = [c["text"] for c in chunks]
    embeddings = model.encode(texts, batch_size=32, show_progress_bar=False)
    for chunk, emb in zip(chunks, embeddings):
        chunk["embedding"] = emb

    save_cached_embeddings(file_key, chunks)
    return chunks


# ============================================================
# AMÉLIORATION 3 — SEMANTIC CHUNKING
# Découpe par paragraphes/sections au lieu d'une taille fixe
# ============================================================

def semantic_chunk(pages_text: dict, max_chunk_size: int = 2000, overlap: int = 200) -> list:
    """
    Chunking sémantique :
    1. Découpe d'abord par paragraphes (blocs naturels)
    2. Fusionne les petits paragraphes jusqu'à max_chunk_size
    3. Conserve le tracking de page
    Avantage : les chunks respectent les frontières de sens du document.
    """
    chunks = []

    for page_num, text in sorted(pages_text.items()):
        # Paragraphes = blocs séparés par lignes vides ou titres (##, numéros)
        paragraphs = re.split(r'\n{2,}|(?=\n[A-Z][^a-z\n]{0,60}\n)', text)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        current_text = ""
        for para in paragraphs:
            if len(current_text) + len(para) + 1 <= max_chunk_size:
                current_text += ("\n\n" if current_text else "") + para
            else:
                if current_text:
                    chunks.append({"text": current_text, "pages": [page_num]})
                # Si paragraphe seul > max_chunk_size → découpe mécanique
                if len(para) > max_chunk_size:
                    for i in range(0, len(para), max_chunk_size - overlap):
                        sub = para[i:i + max_chunk_size].strip()
                        if sub:
                            chunks.append({"text": sub, "pages": [page_num]})
                    current_text = para[-(overlap):] if len(para) > overlap else para
                else:
                    current_text = para

        if current_text:
            chunks.append({"text": current_text, "pages": [page_num]})

    return chunks


def split_into_chunks(pages_text: dict, chunk_size: int = 2000, overlap: int = 200) -> tuple:
    """
    Pipeline de chunking :
    - Essaie le semantic chunking en priorité
    - Conserve le texte complet pour les docs courts
    """
    chunks = semantic_chunk(pages_text, max_chunk_size=chunk_size, overlap=overlap)

    # Reconstruit le texte complet pour les docs courts (mode non-RAG)
    full_text = "\n".join(
        text for _, text in sorted(pages_text.items())
    )

    return chunks, full_text


# ============================================================
# BM25 SIMPLIFIÉ (conservé pour hybrid search)
# ============================================================

def bm25_score(chunk: dict, question: str) -> float:
    q_words = set(re.findall(r'\b\w{3,}\b', question.lower()))
    c_words = set(re.findall(r'\b\w{3,}\b', chunk["text"].lower()))
    if not q_words:
        return 0.0
    return len(q_words & c_words) / len(q_words)


# ============================================================
# HYBRID RETRIEVAL (BM25 + Embeddings + RRF)
# ============================================================

def rrf_score(rank: int, k: int = 60) -> float:
    return 1.0 / (k + rank + 1)


def retrieve_hybrid(chunks: list, question: str, top_k: int = 4, model=None) -> tuple:
    n = len(chunks)

    # Classement BM25
    bm25_scores = [(bm25_score(c, question), i) for i, c in enumerate(chunks)]
    bm25_ranked = [i for _, i in sorted(bm25_scores, key=lambda x: -x[0])]

    # Classement sémantique
    # AMÉLIORATION 5 — fix cohérence : encode([question])[0] au lieu de encode(question)
    if model is not None and all("embedding" in c for c in chunks):
        query_emb = model.encode([question])[0]  # ← fix recommandé
        sem_scores = [
            (cosine_similarity(query_emb, c["embedding"]), i)
            for i, c in enumerate(chunks)
        ]
        sem_ranked = [i for _, i in sorted(sem_scores, key=lambda x: -x[0])]
    else:
        sem_ranked = bm25_ranked  # fallback BM25

    # Fusion RRF
    rrf = {}
    for rank, idx in enumerate(bm25_ranked):
        rrf[idx] = rrf.get(idx, 0.0) + rrf_score(rank)
    for rank, idx in enumerate(sem_ranked):
        rrf[idx] = rrf.get(idx, 0.0) + rrf_score(rank)

    top_indices = sorted(
        sorted(rrf.items(), key=lambda x: -x[1])[:top_k],
        key=lambda x: x[0]
    )

    selected = [chunks[i] for i, _ in top_indices]
    context = "\n\n---\n\n".join(c["text"] for c in selected)

    all_pages = []
    for c in selected:
        all_pages.extend(c["pages"])
    source_pages = sorted(set(all_pages))

    return context, source_pages, selected


# ============================================================
# RERANKING (Cross-Encoder)
# ============================================================

def rerank_chunks(chunks_selected: list, question: str, reranker, top_k: int = 3) -> list:
    if reranker is None or not chunks_selected:
        return chunks_selected
    pairs = [(question, c["text"]) for c in chunks_selected]
    scores = reranker.predict(pairs)
    reranked = sorted(zip(scores, chunks_selected), key=lambda x: -x[0])
    return [c for _, c in reranked[:top_k]]


# ============================================================
# AMÉLIORATION 4 — VECTOR DB FAISS (persistante + scalable)
# Remplace la recherche linéaire en mémoire
# ============================================================

FAISS_INDEX_PATH = os.path.join(CACHE_DIR, "faiss_index.pkl")


def build_faiss_index(chunks: list, file_key: str):
    """
    Construit un index FAISS à partir des embeddings des chunks.
    Sauvegarde l'index sur disque pour éviter la reconstruction.
    """
    try:
        import faiss

        cache_key = f"faiss_{hashlib.md5(file_key.encode()).hexdigest()}"
        faiss_path = os.path.join(CACHE_DIR, f"{cache_key}.faiss")
        meta_path = os.path.join(CACHE_DIR, f"{cache_key}_meta.pkl")

        # Charge index existant
        if os.path.exists(faiss_path) and os.path.exists(meta_path):
            index = faiss.read_index(faiss_path)
            with open(meta_path, "rb") as f:
                metadata = pickle.load(f)
            return index, metadata

        # Construit l'index
        embeddings = np.array([c["embedding"] for c in chunks if "embedding" in c], dtype=np.float32)
        if len(embeddings) == 0:
            return None, None

        dim = embeddings.shape[1]
        # IndexFlatIP + normalisation = cosine similarity via produit interne
        faiss.normalize_L2(embeddings)
        index = faiss.IndexFlatIP(dim)
        index.add(embeddings)

        # Sauvegarde
        faiss.write_index(index, faiss_path)
        metadata = [{"text": c["text"], "pages": c["pages"]} for c in chunks if "embedding" in c]
        with open(meta_path, "wb") as f:
            pickle.dump(metadata, f)

        return index, metadata

    except ImportError:
        return None, None  # Fallback linéaire si FAISS non installé


def faiss_search(index, metadata: list, query_emb: np.ndarray, top_k: int) -> list:
    """Recherche dans l'index FAISS. Retourne les top_k chunks les plus proches."""
    try:
        import faiss
        query = np.array([query_emb], dtype=np.float32)
        faiss.normalize_L2(query)
        _, indices = index.search(query, top_k)
        return [metadata[i] for i in indices[0] if i < len(metadata)]
    except Exception:
        return []


def retrieve_hybrid_faiss(chunks: list, question: str, top_k: int = 6, model=None, file_key: str = "") -> tuple:
    """
    Hybrid retrieval avec FAISS si disponible, sinon fallback linéaire.
    """
    # Tentative FAISS
    if model is not None and file_key:
        faiss_index, faiss_meta = build_faiss_index(chunks, file_key)
        if faiss_index is not None and faiss_meta is not None:
            query_emb = model.encode([question])[0]

            # Recherche FAISS (sémantique)
            faiss_results = faiss_search(faiss_index, faiss_meta, query_emb, top_k)

            # Recherche BM25
            bm25_scores_list = [(bm25_score(c, question), i) for i, c in enumerate(chunks)]
            bm25_top = [chunks[i] for _, i in sorted(bm25_scores_list, key=lambda x: -x[0])[:top_k]]

            # Fusion simple : union des deux listes (dédupliquée par texte)
            seen = set()
            combined = []
            for c in faiss_results + bm25_top:
                key = c["text"][:100]
                if key not in seen:
                    seen.add(key)
                    combined.append(c)

            selected = combined[:top_k]
            context = "\n\n---\n\n".join(c["text"] for c in selected)
            all_pages = []
            for c in selected:
                all_pages.extend(c["pages"])
            return context, sorted(set(all_pages)), selected

    # Fallback : hybrid retrieval linéaire (original)
    return retrieve_hybrid(chunks, question, top_k=top_k, model=model)


# ============================================================
# PIPELINE RETRIEVAL COMPLET
# ============================================================

def ask_full_or_rag(client, question: str) -> tuple:
    full_text = st.session_state.get("full_text", "")
    chunks = st.session_state.get("chunks", [])
    file_key = st.session_state.get("loaded_file", "")

    if not full_text:
        return "Aucun document chargé.", []

    embedding_model = load_embedding_model()
    reranker = load_reranker()

    if len(full_text) <= 25000:
        response = ask_mistral(client, full_text, question)
        _, source_pages, _ = retrieve_hybrid_faiss(
            chunks, question, top_k=3, model=embedding_model, file_key=file_key
        )
        return response, source_pages
    else:
        # AMÉLIORATION 6 — top_k_retrieve par défaut = 10 (plus large)
        top_k_retrieve = st.session_state.get("top_k", 10)
        top_k_rerank = st.session_state.get("top_k_rerank", 3)

        # Étape 1 : Hybrid retrieval (avec FAISS si dispo)
        context_raw, _, candidates = retrieve_hybrid_faiss(
            chunks, question, top_k=top_k_retrieve, model=embedding_model, file_key=file_key
        )

        # Étape 2 : Reranking cross-encoder
        reranked = rerank_chunks(candidates, question, reranker, top_k=top_k_rerank)

        context = "\n\n---\n\n".join(c["text"] for c in reranked)
        all_pages = []
        for c in reranked:
            all_pages.extend(c["pages"])
        source_pages = sorted(set(all_pages))

        response = ask_mistral(client, context, question)
        return response, source_pages


# ============================================================
# CLIENT MISTRAL
# ============================================================

def get_client():
    api_key = st.secrets.get("MISTRAL_API_KEY") or os.getenv("MISTRAL_API_KEY")
    if not api_key:
        return None
    return Mistral(api_key=api_key)


def ask_mistral(client, context: str, question: str) -> str:
    try:
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": f"CONTEXTE:\n{context}\n\nQUESTION: {question}"}
            ],
            temperature=0,
            max_tokens=1500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Erreur Mistral : {e}"


def format_sources(pages: list) -> str:
    if not pages:
        return ""
    if len(pages) == 1:
        return f"📄 Source : Page {pages[0]}"
    return f"📄 Sources : Pages {', '.join(str(p) for p in pages)}"


# ============================================================
# AMÉLIORATION 7 — VRAIE ÉVALUATION RAGAS
# Utilise la lib ragas si installée, sinon fallback LLM-as-judge
# ============================================================

def evaluate_with_ragas(question: str, answer: str, contexts: list) -> dict:
    """
    Évaluation avec la vraie lib RAGAS (ragas.evaluate).
    Retourne None si ragas n'est pas installé.
    """
    try:
        from ragas import evaluate
        from ragas.metrics import faithfulness, answer_relevancy, context_recall
        from datasets import Dataset

        dataset = Dataset.from_dict({
            "question": [question],
            "answer": [answer],
            "contexts": [contexts],
        })

        result = evaluate(
            dataset=dataset,
            metrics=[faithfulness, answer_relevancy, context_recall],
        )

        return {
            "faithfulness": float(result["faithfulness"]),
            "answer_relevance": float(result["answer_relevancy"]),
            "context_recall": float(result["context_recall"]),
            "faithfulness_reason": "Calculé via ragas.evaluate()",
            "answer_relevance_reason": "Calculé via ragas.evaluate()",
            "context_recall_reason": "Calculé via ragas.evaluate()",
            "source": "ragas"
        }
    except ImportError:
        return None  # Fallback LLM-as-judge
    except Exception as e:
        return {"error": str(e)}


def evaluate_rag_answer(client, question: str, context: str, answer: str, chunks_selected: list = None) -> dict:
    """
    Évaluation RAG :
    1. Essaie d'abord la vraie lib RAGAS
    2. Fallback sur LLM-as-judge (approximation)
    """
    # Tentative vraie RAGAS
    if chunks_selected:
        contexts_list = [c["text"] for c in chunks_selected]
        ragas_result = evaluate_with_ragas(question, answer, contexts_list)
        if ragas_result is not None:
            return ragas_result

    # Fallback LLM-as-judge
    eval_prompt = f"""Tu es un évaluateur RAG expert. Évalue les 3 métriques suivantes
en retournant UNIQUEMENT un JSON valide, sans markdown :

QUESTION: {question}

CONTEXTE UTILISÉ:
{context[:3000]}

RÉPONSE GÉNÉRÉE:
{answer}

Réponds avec ce format exact :
{{
  "faithfulness": <float 0.0-1.0>,
  "answer_relevance": <float 0.0-1.0>,
  "context_recall": <float 0.0-1.0>,
  "faithfulness_reason": "<explication courte>",
  "answer_relevance_reason": "<explication courte>",
  "context_recall_reason": "<explication courte>"
}}

Définitions :
- faithfulness : la réponse ne contient que des infos du contexte (1.0 = totalement fidèle)
- answer_relevance : la réponse répond précisément à la question (1.0 = parfait)
- context_recall : le contexte contient les infos pour répondre (1.0 = contexte complet)"""

    try:
        raw = client.chat.complete(
            model="mistral-large-latest",
            messages=[{"role": "user", "content": eval_prompt}],
            temperature=0,
            max_tokens=500
        ).choices[0].message.content.strip()

        raw = re.sub(r"```(?:json)?", "", raw).strip("` \n")
        start, end = raw.find("{"), raw.rfind("}") + 1
        result = json.loads(raw[start:end])
        result["source"] = "llm-as-judge"
        return result
    except Exception as e:
        return {"error": str(e)}


def render_metric_bar(label: str, value: float, reason: str):
    color = "#4caf50" if value >= 0.7 else "#ff9800" if value >= 0.4 else "#f44336"
    st.markdown(f"""
    <div style="margin-bottom:12px">
        <div style="display:flex;justify-content:space-between;margin-bottom:4px">
            <span style="font-weight:600">{label}</span>
            <span style="color:{color};font-weight:700">{value:.2f}</span>
        </div>
        <div style="background:#e0e0e0;border-radius:8px;height:10px">
            <div style="background:{color};width:{int(value*100)}%;height:10px;border-radius:8px"></div>
        </div>
        <div style="font-size:0.8em;color:#666;margin-top:4px">{reason}</div>
    </div>
    """, unsafe_allow_html=True)


# ============================================================
# EXPORT PDF CONVERSATION
# ============================================================

def export_chat_to_pdf(messages: list, doc_name: str) -> bytes:
    def clean(text: str) -> str:
        return text.encode("latin-1", errors="replace").decode("latin-1")

    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    pdf.set_font("Helvetica", "B", 16)
    pdf.set_fill_color(30, 60, 114)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 12, clean("Conversation - Insight PDF Pro"), fill=True, ln=True, align="C")
    pdf.ln(2)

    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 8, clean(f"Document : {doc_name}"), ln=True, align="C")
    pdf.ln(6)

    for msg in messages:
        role = msg["role"]
        text = clean(msg["content"])
        pages = msg.get("pages", [])

        if role == "user":
            pdf.set_fill_color(230, 240, 255)
            pdf.set_text_color(30, 60, 114)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 8, "Vous :", ln=True, fill=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 30, 30)
            pdf.multi_cell(0, 7, text)
        else:
            pdf.set_fill_color(245, 245, 245)
            pdf.set_text_color(80, 80, 80)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 8, "Assistant :", ln=True, fill=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 30, 30)
            pdf.multi_cell(0, 7, text)
            if pages:
                pdf.set_font("Helvetica", "I", 9)
                pdf.set_text_color(100, 100, 200)
                src = (
                    f"Sources : Pages {', '.join(str(p) for p in pages)}"
                    if len(pages) > 1 else f"Source : Page {pages[0]}"
                )
                pdf.cell(0, 6, clean(src), ln=True)

        pdf.ln(4)

    return bytes(pdf.output())


# ============================================================
# UTILITAIRES
# ============================================================

def create_pptx(data):
    prs = Presentation()
    bg_color = RGBColor(30, 60, 114)
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tb.text_frame
        tf.text = slide_data.get("titre", "Slide")
        p = tf.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(24)

        content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        cf = content.text_frame
        cf.word_wrap = True
        for pt in slide_data.get("points", []):
            para = cf.add_paragraph()
            para.text = f"• {pt}"
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(255, 255, 255)

    ppt_io = BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()


# ============================================================
# INTERFACE
# ============================================================

st.title("✨ Insight PDF Pro")
st.caption(f"Développé par Herman Kandolo • {datetime.now().year}")

client = get_client()
if not client:
    st.error("⚠️ Clé API Mistral manquante. Ajoutez MISTRAL_API_KEY dans les secrets Streamlit.")
    st.stop()

# --- SIDEBAR ---
with st.sidebar:
    st.subheader("📤 Importation")
    uploaded_file = st.file_uploader("Choisir un PDF", type="pdf", label_visibility="collapsed")

    if uploaded_file:
        file_key = uploaded_file.name
        if st.session_state.get("loaded_file") != file_key:
            with st.spinner("Extraction et découpage du texte…"):
                pages_text = extract_pdf_data(uploaded_file)
                if not pages_text:
                    st.error("Le PDF semble vide ou non lisible (PDF scanné ?).")
                    st.stop()
                # Semantic chunking
                chunks, full_text = split_into_chunks(pages_text)

                # Embeddings avec cache persistant
                emb_model = load_embedding_model()
                if emb_model:
                    with st.spinner("Encodage sémantique des chunks (avec cache)…"):
                        chunks = encode_chunks(chunks, emb_model, file_key)

                st.session_state.pdf_pages = pages_text
                st.session_state.full_text = full_text
                st.session_state.chunks = chunks
                st.session_state.messages = []
                st.session_state.loaded_file = file_key

            # Status des améliorations actives
            emb_status = "✅ embeddings" if emb_model else "⚠️ BM25 only"
            try:
                import faiss
                faiss_status = "✅ FAISS"
            except ImportError:
                faiss_status = "⚠️ no FAISS"
            try:
                import pdfplumber
                parser_status = "pdfplumber"
            except ImportError:
                parser_status = "PyPDF2"

            st.success(
                f"✅ {len(pages_text)} pages • {len(chunks)} chunks\n"
                f"{emb_status} • {faiss_status} • parser: {parser_status}"
            )
        else:
            st.info(f"📄 {file_key} déjà chargé.")

    if "pdf_pages" in st.session_state:
        with st.expander("ℹ️ Détails & Paramètres RAG"):
            st.metric("Pages", len(st.session_state.pdf_pages))
            st.metric("Chunks RAG", len(st.session_state.get("chunks", [])))
            st.metric("Caractères", f"{len(st.session_state.full_text):,}")
            st.divider()
            st.caption("🔍 Retrieval")
            # AMÉLIORATION 6 — défaut = 10 (plus large pour meilleur recall)
            st.session_state.top_k = st.slider(
                "Chunks candidats (retrieval)", 3, 12,
                st.session_state.get("top_k", 10),
                help="Nombre de chunks récupérés avant reranking (recommandé : 8-12)"
            )
            st.caption("🏆 Reranking")
            st.session_state.top_k_rerank = st.slider(
                "Chunks finaux (après reranking)", 1, 5,
                st.session_state.get("top_k_rerank", 3),
                help="Chunks envoyés au LLM après cross-encoder (recommandé : 3-5)"
            )


# --- ONGLETS ---
if "pdf_pages" in st.session_state:
    tabs = st.tabs([
        "💬 Chat", "📝 Synthèse", "📊 Analyse",
        "🔊 Audio", "🎯 Présentation", "📐 Évaluation RAG"
    ])

    # ── TAB 1 : CHAT ─────────────────────────────────────────
    with tabs[0]:
        is_long = len(st.session_state.full_text) > 25000
        emb_ready = all("embedding" in c for c in st.session_state.get("chunks", []))
        reranker_ready = load_reranker() is not None

        col_info, col_export = st.columns([4, 1])
        with col_info:
            flags = []
            if is_long:
                flags.append(f"RAG actif ({len(st.session_state.chunks)} chunks)")
            if emb_ready:
                flags.append("embeddings ✅")
            if reranker_ready:
                flags.append("reranking ✅")
            try:
                import faiss
                flags.append("FAISS ✅")
            except ImportError:
                pass
            if flags:
                st.caption("📚 " + " • ".join(flags))

        with col_export:
            if st.session_state.get("messages"):
                pdf_bytes = export_chat_to_pdf(
                    st.session_state.messages,
                    st.session_state.get("loaded_file", "document")
                )
                st.download_button(
                    label="⬇️ PDF",
                    data=pdf_bytes,
                    file_name="conversation.pdf",
                    mime="application/pdf",
                    key="dl_chat_pdf"
                )

        if "messages" not in st.session_state:
            st.session_state.messages = []

        for msg in st.session_state.messages:
            with st.chat_message(msg["role"]):
                st.write(msg["content"])
                if msg["role"] == "assistant" and msg.get("pages"):
                    st.caption(format_sources(msg["pages"]))

        if prompt := st.chat_input("Posez une question sur le document…"):
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)
            with st.chat_message("assistant", avatar="✨"):
                with st.spinner("Recherche dans le document…"):
                    response, source_pages = ask_full_or_rag(client, prompt)
                    st.write(response)
                    if source_pages:
                        st.caption(format_sources(source_pages))
                    st.session_state.messages.append({
                        "role": "assistant",
                        "content": response,
                        "pages": source_pages
                    })

    # ── TAB 2 : SYNTHÈSE ────────────────────────────────────
    with tabs[1]:
        s_mode = st.select_slider(
            "Niveau de précision",
            options=["Court", "Moyen", "Détaillé"],
            value="Moyen"
        )
        longueur = {
            "Court": "en 5 phrases",
            "Moyen": "en 10-15 phrases",
            "Détaillé": "de manière exhaustive"
        }

        if st.button("📝 Rédiger le résumé", key="btn_resume"):
            with st.spinner("Génération du résumé…"):
                full_text = st.session_state.full_text
                if len(full_text) > 25000:
                    chunks = st.session_state.chunks
                    step = max(1, len(chunks) // 8)
                    sampled = chunks[::step][:8]
                    context = "\n\n---\n\n".join(c["text"] for c in sampled)
                    all_pages = []
                    for c in sampled:
                        all_pages.extend(c["pages"])
                    source_pages = sorted(set(all_pages))
                else:
                    context = full_text
                    source_pages = sorted(st.session_state.pdf_pages.keys())

                question = (
                    f"Fais un résumé structuré {longueur[s_mode]} de ce document, "
                    f"avec des sections claires."
                )
                result = ask_mistral(client, context, question)
                st.info(result)
                st.caption(format_sources(source_pages))

    # ── TAB 3 : ANALYSE ─────────────────────────────────────
    with tabs[2]:
        col1, col2 = st.columns(2)
        words = re.findall(r'\b\w+\b', st.session_state.full_text.lower())
        stop_words = {
            "les", "des", "une", "que", "qui", "dans", "pour", "avec", "sur",
            "par", "est", "sont", "this", "that", "from", "have", "been",
            "will", "leur", "leurs", "mais", "donc", "comme", "plus", "aussi",
            "tout", "tous", "très", "bien", "être", "avoir", "faire"
        }
        freq = Counter([w for w in words if len(w) > 3 and w not in stop_words])

        with col1:
            st.metric("Mots totaux", f"{len(words):,}")
            st.metric("Pages analysées", len(st.session_state.pdf_pages))
            st.metric("Chunks créés", len(st.session_state.get("chunks", [])))
            st.subheader("🔑 Mots-clés fréquents")
            for w, c in freq.most_common(10):
                st.write(f"- **{w}** : {c} occurrences")

        with col2:
            if st.button("🔍 Analyse sémantique", key="btn_semantic"):
                with st.spinner("Analyse en cours…"):
                    question = (
                        "Quels sont les thèmes principaux de ce document ? "
                        "Liste-les et explique chacun brièvement."
                    )
                    result, source_pages = ask_full_or_rag(client, question)
                    st.write(result)
                    st.caption(format_sources(source_pages))

    # ── TAB 4 : AUDIO ───────────────────────────────────────
    with tabs[3]:
        max_page = len(st.session_state.pdf_pages)
        p_num = st.number_input("Numéro de page à lire", min_value=1, max_value=max_page, value=1)
        lang = st.selectbox("Langue", ["fr", "en", "es", "de"], index=0)

        if st.button("🔊 Générer l'audio", key="btn_audio"):
            page_text = st.session_state.pdf_pages.get(p_num, "")
            if page_text:
                with st.spinner("Génération audio…"):
                    try:
                        tts = gTTS(text=page_text, lang=lang)
                        audio_io = BytesIO()
                        tts.write_to_fp(audio_io)
                        audio_io.seek(0)
                        st.audio(audio_io, format="audio/mp3")
                        st.caption(f"📄 Page {p_num}")
                    except Exception as e:
                        st.error(f"Erreur audio : {e}")
            else:
                st.warning("Aucun texte trouvé sur cette page.")

    # ── TAB 5 : PRÉSENTATION ────────────────────────────────
    with tabs[4]:
        n_slides = st.number_input("Nombre de slides", min_value=3, max_value=10, value=5)

        if st.button("🎯 Générer la présentation PPTX", key="btn_pptx"):
            with st.spinner("L'IA structure vos slides…"):
                try:
                    question = (
                        f"Crée une structure pour exactement {n_slides} slides basées sur ce document. "
                        f"Réponds UNIQUEMENT avec un JSON valide, sans texte avant ou après, "
                        f'sans balises markdown : '
                        f'{{ "slides": [ {{ "titre": "Titre de la slide", '
                        f'"points": ["Point 1", "Point 2", "Point 3"] }} ] }}'
                    )
                    raw, source_pages = ask_full_or_rag(client, question)
                    raw = raw.strip()
                    if raw.startswith("```"):
                        raw = re.sub(r"```(?:json)?", "", raw).strip("` \n")

                    start = raw.find('{')
                    end = raw.rfind('}') + 1

                    if start != -1 and end > 0:
                        data = json.loads(raw[start:end])
                        ppt_bytes = create_pptx(data)
                        st.download_button(
                            label="📥 Télécharger la présentation",
                            data=ppt_bytes,
                            file_name="presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.success(f"✅ {len(data.get('slides', []))} slides générées !")
                        st.caption(format_sources(source_pages))
                    else:
                        st.error("Format JSON invalide reçu.")
                        st.code(raw)
                except json.JSONDecodeError as e:
                    st.error(f"Erreur JSON : {e}")
                    st.code(raw)
                except Exception as e:
                    st.error(f"Erreur : {e}")

    # ── TAB 6 : ÉVALUATION RAG ──────────────────────────────
    with tabs[5]:
        st.subheader("📐 Évaluation RAG")

        # Indique la méthode d'évaluation disponible
        try:
            import ragas
            st.success("✅ Vraie lib **RAGAS** détectée — évaluation de haute fidélité activée.")
        except ImportError:
            st.info("ℹ️ RAGAS non installé. Utilisation du mode LLM-as-judge (approximation). "
                    "`pip install ragas datasets` pour activer l'évaluation exacte.")

        st.caption(
            "Évalue la qualité de ton pipeline RAG sur 3 métriques : "
            "Faithfulness · Answer Relevance · Context Recall"
        )

        last_q, last_a = "", ""
        if st.session_state.get("messages"):
            msgs = st.session_state.messages
            for i in range(len(msgs) - 1, -1, -1):
                if msgs[i]["role"] == "assistant" and i > 0:
                    last_a = msgs[i]["content"]
                    last_q = msgs[i - 1]["content"]
                    break

        eval_q = st.text_area("Question à évaluer", value=last_q, height=80)
        eval_a = st.text_area("Réponse à évaluer", value=last_a, height=120)

        if st.button("📊 Lancer l'évaluation", key="btn_eval"):
            if not eval_q or not eval_a:
                st.warning("Renseigne une question et une réponse.")
            else:
                with st.spinner("Évaluation en cours…"):
                    emb_model = load_embedding_model()
                    file_key = st.session_state.get("loaded_file", "")
                    context, source_pages, chunks_selected = retrieve_hybrid_faiss(
                        st.session_state.chunks, eval_q,
                        top_k=st.session_state.get("top_k", 10),
                        model=emb_model,
                        file_key=file_key
                    )
                    metrics = evaluate_rag_answer(
                        client, eval_q, context, eval_a,
                        chunks_selected=chunks_selected
                    )

                if "error" in metrics:
                    st.error(f"Erreur évaluation : {metrics['error']}")
                else:
                    source_label = metrics.get("source", "llm-as-judge")
                    st.caption(f"🔬 Méthode : **{source_label}**")

                    st.markdown("### Résultats")
                    col_m1, col_m2, col_m3 = st.columns(3)
                    with col_m1:
                        st.metric("Faithfulness", f"{metrics.get('faithfulness', 0):.2f}")
                    with col_m2:
                        st.metric("Answer Relevance", f"{metrics.get('answer_relevance', 0):.2f}")
                    with col_m3:
                        st.metric("Context Recall", f"{metrics.get('context_recall', 0):.2f}")

                    st.divider()
                    render_metric_bar(
                        "Faithfulness",
                        metrics.get("faithfulness", 0),
                        metrics.get("faithfulness_reason", "")
                    )
                    render_metric_bar(
                        "Answer Relevance",
                        metrics.get("answer_relevance", 0),
                        metrics.get("answer_relevance_reason", "")
                    )
                    render_metric_bar(
                        "Context Recall",
                        metrics.get("context_recall", 0),
                        metrics.get("context_recall_reason", "")
                    )

                    avg = np.mean([
                        metrics.get("faithfulness", 0),
                        metrics.get("answer_relevance", 0),
                        metrics.get("context_recall", 0)
                    ])
                    color = "#4caf50" if avg >= 0.7 else "#ff9800" if avg >= 0.4 else "#f44336"
                    st.markdown(
                        f"<h3 style='color:{color}'>Score global : {avg:.2f} / 1.00</h3>",
                        unsafe_allow_html=True
                    )
                    st.caption(format_sources(source_pages))

                    if "eval_history" not in st.session_state:
                        st.session_state.eval_history = []
                    st.session_state.eval_history.append({
                        "question": eval_q[:60] + "…",
                        "faithfulness": metrics.get("faithfulness", 0),
                        "answer_relevance": metrics.get("answer_relevance", 0),
                        "context_recall": metrics.get("context_recall", 0),
                        "avg": avg,
                        "méthode": source_label
                    })

        if st.session_state.get("eval_history"):
            st.divider()
            st.subheader("📋 Historique des évaluations")
            import pandas as pd
            df = pd.DataFrame(st.session_state.eval_history)
            st.dataframe(df, use_container_width=True)

else:
    st.info("👈 Veuillez charger un fichier PDF dans la barre latérale pour commencer.")
