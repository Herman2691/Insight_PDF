import streamlit as st
import PyPDF2
from io import BytesIO
import json
import re
from collections import Counter
from datetime import datetime
import os
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from mistralai import Mistral
from fpdf import FPDF

# --- CONFIGURATION ---
st.set_page_config(page_title="Insight PDF Pro", page_icon="✨", layout="wide")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Google+Sans:wght@400;500;700&display=swap');
html, body, [class*="css"] { font-family: 'Google Sans', sans-serif; }
.stButton>button {
    border-radius: 20px;
    background-color: #f0f4f9;
    border: none;
    color: #1a73e8;
    font-weight: 500;
}
[data-testid="stSidebar"] {
    background-color: #f8f9fa;
    border-right: 1px solid #e1e3e1;
}
/* Bulles messages */
[data-testid="stChatMessage"] {
    border-radius: 16px;
    margin-bottom: 8px;
    padding: 4px 8px;
}
</style>
""", unsafe_allow_html=True)

# ============================================================
# PROMPT ANTI-HALLUCINATION
# ============================================================

SYSTEM_PROMPT = (
    "Tu es un assistant expert en analyse de documents. "
    "Réponds UNIQUEMENT en utilisant les informations du CONTEXTE fourni. "
    "Si la réponse ne se trouve pas dans le contexte, réponds exactement : "
    "'Je suis désolé, mais cette information n'est pas présente dans le document fourni.' "
    "Ne réponds jamais en utilisant tes connaissances générales si le sujet est absent du document."
)

# ============================================================
# CHUNKING MANUEL AVEC NUMÉROS DE PAGE
# ============================================================

def split_into_chunks(pages_text: dict, chunk_size: int = 2000, overlap: int = 200) -> list:
    """
    Découpe le texte en chunks en conservant le numéro de page source.
    Retourne une liste de dicts : {"text": "...", "pages": [3, 4]}
    """
    chunks = []

    # On construit une liste de (page_num, texte) pour tracker la page d'origine
    page_boundaries = []
    full_text = ""
    for page_num, text in sorted(pages_text.items()):
        start_pos = len(full_text)
        full_text += text + "\n"
        end_pos = len(full_text)
        page_boundaries.append((page_num, start_pos, end_pos))

    # Découpe en chunks avec chevauchement
    start = 0
    text_len = len(full_text)

    while start < text_len:
        end = start + chunk_size

        if end < text_len:
            cut = max(
                full_text.rfind(". ", start, end),
                full_text.rfind("\n", start, end)
            )
            if cut > start + chunk_size // 2:
                end = cut + 1

        chunk_text = full_text[start:end].strip()

        if chunk_text:
            # Trouve quelles pages sont couvertes par ce chunk
            covered_pages = []
            for page_num, p_start, p_end in page_boundaries:
                if p_start < end and p_end > start:
                    covered_pages.append(page_num)

            chunks.append({
                "text": chunk_text,
                "pages": covered_pages
            })

        start = end - overlap

    return chunks, full_text


def score_chunk(chunk: dict, question: str) -> float:
    q_words = set(re.findall(r'\b\w{3,}\b', question.lower()))
    c_words = set(re.findall(r'\b\w{3,}\b', chunk["text"].lower()))
    if not q_words:
        return 0.0
    return len(q_words & c_words) / len(q_words)


def retrieve_best_chunks(chunks: list, question: str, top_k: int = 4) -> tuple:
    """
    Retourne (contexte_texte, pages_sources).
    """
    scored = [(score_chunk(c, question), i, c) for i, c in enumerate(chunks)]
    scored.sort(key=lambda x: (-x[0], x[1]))
    best = scored[:top_k]
    best.sort(key=lambda x: x[1])  # retri par ordre naturel

    context = "\n\n---\n\n".join(c["text"] for _, _, c in best)

    # Pages sources uniques triées
    all_pages = []
    for _, _, c in best:
        all_pages.extend(c["pages"])
    source_pages = sorted(set(all_pages))

    return context, source_pages


# ============================================================
# CLIENT MISTRAL
# ============================================================

def get_client():
    api_key = st.secrets.get("MISTRAL_API_KEY") or os.getenv("MISTRAL_API_KEY")
    if not api_key:
        return None
    return Mistral(api_key=api_key)


def ask_mistral(client, context: str, question: str) -> str:
    try:
        response = client.chat.complete(
            model="mistral-large-latest",
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": f"CONTEXTE:\n{context}\n\nQUESTION: {question}"}
            ],
            temperature=0,
            max_tokens=1500
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Erreur Mistral : {e}"


def ask_full_or_rag(client, question: str) -> tuple:
    """
    Retourne (réponse, pages_sources).
    - Doc court  → texte complet, pages = toutes les pages
    - Doc long   → RAG, pages = pages des chunks sélectionnés
    """
    full_text = st.session_state.get("full_text", "")
    chunks = st.session_state.get("chunks", [])

    if not full_text:
        return "Aucun document chargé.", []

    if len(full_text) <= 25000:
        response = ask_mistral(client, full_text, question)
        # Pour un doc court on ne peut pas savoir la page exacte sans RAG
        # On retourne les pages des chunks qui scorent le mieux
        _, source_pages = retrieve_best_chunks(chunks, question, top_k=3)
        return response, source_pages
    else:
        top_k = st.session_state.get("top_k", 4)
        context, source_pages = retrieve_best_chunks(chunks, question, top_k=top_k)
        response = ask_mistral(client, context, question)
        return response, source_pages


def format_sources(pages: list) -> str:
    """Formate l'affichage des pages sources."""
    if not pages:
        return ""
    if len(pages) == 1:
        return f"📄 Source : Page {pages[0]}"
    return f"📄 Sources : Pages {', '.join(str(p) for p in pages)}"


# ============================================================
# EXPORT PDF CONVERSATION
# ============================================================

def export_chat_to_pdf(messages: list, doc_name: str) -> bytes:
    """Génère un PDF de la conversation."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)

    # Titre
    pdf.set_font("Helvetica", "B", 16)
    pdf.set_fill_color(30, 60, 114)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 12, "Conversation - Insight PDF Pro", fill=True, ln=True, align="C")
    pdf.ln(2)

    # Sous-titre
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(120, 120, 120)
    pdf.cell(0, 8, f"Document : {doc_name}", ln=True, align="C")
    pdf.ln(6)

    for msg in messages:
        role = msg["role"]
        content = msg["content"]
        pages = msg.get("pages", [])

        if role == "user":
            # Bulle utilisateur
            pdf.set_fill_color(230, 240, 255)
            pdf.set_text_color(30, 60, 114)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 8, "Vous :", ln=True, fill=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 30, 30)
            pdf.multi_cell(0, 7, content)
        else:
            # Bulle assistant
            pdf.set_fill_color(245, 245, 245)
            pdf.set_text_color(80, 80, 80)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 8, "Assistant :", ln=True, fill=True)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(30, 30, 30)
            pdf.multi_cell(0, 7, content)
            if pages:
                pdf.set_font("Helvetica", "I", 9)
                pdf.set_text_color(100, 100, 200)
                src = f"Sources : Pages {', '.join(str(p) for p in pages)}" if len(pages) > 1 else f"Source : Page {pages[0]}"
                pdf.cell(0, 6, src, ln=True)

        pdf.ln(4)

    return bytes(pdf.output())


# ============================================================
# FONCTIONS UTILITAIRES
# ============================================================

def extract_pdf_data(pdf_file):
    reader = PyPDF2.PdfReader(BytesIO(pdf_file.read()))
    pages_text = {}
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages_text[i + 1] = text
    return pages_text


def create_pptx(data):
    prs = Presentation()
    bg_color = RGBColor(30, 60, 114)
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = tb.text_frame
        tf.text = slide_data.get("titre", "Slide")
        p = tf.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(24)

        content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        cf = content.text_frame
        cf.word_wrap = True
        for pt in slide_data.get("points", []):
            para = cf.add_paragraph()
            para.text = f"• {pt}"
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(255, 255, 255)

    ppt_io = BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()


# ============================================================
# INTERFACE
# ============================================================

st.title("✨ Insight PDF Pro")
st.caption(f"Développé par Herman Kandolo • {datetime.now().year}")

client = get_client()
if not client:
    st.error("⚠️ Clé API Mistral manquante. Ajoutez MISTRAL_API_KEY dans les secrets Streamlit.")
    st.stop()

# --- SIDEBAR ---
with st.sidebar:
    st.subheader("📤 Importation")
    uploaded_file = st.file_uploader("Choisir un PDF", type="pdf", label_visibility="collapsed")

    if uploaded_file:
        file_key = uploaded_file.name
        if st.session_state.get("loaded_file") != file_key:
            with st.spinner("Extraction et découpage du texte..."):
                pages_text = extract_pdf_data(uploaded_file)
                if not pages_text:
                    st.error("Le PDF semble vide ou non lisible (PDF scanné ?).")
                    st.stop()
                chunks, full_text = split_into_chunks(pages_text)
                st.session_state.pdf_pages = pages_text
                st.session_state.full_text = full_text
                st.session_state.chunks = chunks
                st.session_state.messages = []
                st.session_state.loaded_file = file_key
            st.success(f"✅ {len(pages_text)} pages • {len(chunks)} chunks")
        else:
            st.info(f"📄 {file_key} déjà chargé.")

    if "pdf_pages" in st.session_state:
        with st.expander("ℹ️ Détails du document"):
            st.metric("Pages", len(st.session_state.pdf_pages))
            st.metric("Chunks RAG", len(st.session_state.get("chunks", [])))
            st.metric("Caractères", f"{len(st.session_state.full_text):,}")
            st.divider()
            st.session_state.top_k = st.slider(
                "Chunks retenus par requête", 1, 8,
                st.session_state.get("top_k", 4)
            )


# --- ONGLETS ---
if "pdf_pages" in st.session_state:
    tabs = st.tabs(["💬 Chat", "📝 Synthèse", "📊 Analyse", "🔊 Audio", "🎯 Présentation"])

    # TAB 1 : CHAT
    with tabs[0]:
        is_long = len(st.session_state.full_text) > 25000

        # Barre supérieure : info RAG + bouton export
        col_info, col_export = st.columns([4, 1])
        with col_info:
            if is_long:
                st.caption(f"📚 Document long — mode RAG actif ({len(st.session_state.chunks)} chunks)")
        with col_export:
            if st.session_state.get("messages"):
                pdf_bytes = export_chat_to_pdf(
                    st.session_state.messages,
                    st.session_state.get("loaded_file", "document")
                )
                st.download_button(
                    label="⬇️ PDF",
                    data=pdf_bytes,
                    file_name="conversation.pdf",
                    mime="application/pdf",
                    key="dl_chat_pdf"
                )

        if "messages" not in st.session_state:
            st.session_state.messages = []

        # Affichage historique
        for msg in st.session_state.messages:
            with st.chat_message(msg["role"]):
                st.write(msg["content"])
                if msg["role"] == "assistant" and msg.get("pages"):
                    st.caption(format_sources(msg["pages"]))

        if prompt := st.chat_input("Posez une question sur le document..."):
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.write(prompt)
            with st.chat_message("assistant", avatar="✨"):
                with st.spinner("Recherche dans le document..."):
                    response, source_pages = ask_full_or_rag(client, prompt)
                    st.write(response)
                    if source_pages:
                        st.caption(format_sources(source_pages))
                    st.session_state.messages.append({
                        "role": "assistant",
                        "content": response,
                        "pages": source_pages
                    })

    # TAB 2 : SYNTHÈSE
    with tabs[1]:
        s_mode = st.select_slider(
            "Niveau de précision",
            options=["Court", "Moyen", "Détaillé"],
            value="Moyen"
        )
        longueur = {
            "Court": "en 5 phrases",
            "Moyen": "en 10-15 phrases",
            "Détaillé": "de manière exhaustive"
        }

        if st.button("📝 Rédiger le résumé", key="btn_resume"):
            with st.spinner("Génération du résumé..."):
                full_text = st.session_state.full_text
                if len(full_text) > 25000:
                    chunks = st.session_state.chunks
                    step = max(1, len(chunks) // 8)
                    sampled = chunks[::step][:8]
                    context = "\n\n---\n\n".join(c["text"] for c in sampled)
                    all_pages = []
                    for c in sampled:
                        all_pages.extend(c["pages"])
                    source_pages = sorted(set(all_pages))
                else:
                    context = full_text
                    source_pages = sorted(st.session_state.pdf_pages.keys())

                question = f"Fais un résumé structuré {longueur[s_mode]} de ce document, avec des sections claires."
                result = ask_mistral(client, context, question)
                st.info(result)
                st.caption(format_sources(source_pages))

    # TAB 3 : ANALYSE
    with tabs[2]:
        col1, col2 = st.columns(2)
        words = re.findall(r'\b\w+\b', st.session_state.full_text.lower())
        stop_words = {
            "les", "des", "une", "que", "qui", "dans", "pour", "avec", "sur",
            "par", "est", "sont", "this", "that", "from", "have", "been",
            "will", "leur", "leurs", "mais", "donc", "comme", "plus", "aussi",
            "tout", "tous", "très", "bien", "être", "avoir", "faire"
        }
        freq = Counter([w for w in words if len(w) > 3 and w not in stop_words])

        with col1:
            st.metric("Mots totaux", f"{len(words):,}")
            st.metric("Pages analysées", len(st.session_state.pdf_pages))
            st.metric("Chunks créés", len(st.session_state.get("chunks", [])))
            st.subheader("🔑 Mots-clés fréquents")
            for w, c in freq.most_common(10):
                st.write(f"- **{w}** : {c} occurrences")

        with col2:
            if st.button("🔍 Analyse sémantique", key="btn_semantic"):
                with st.spinner("Analyse en cours..."):
                    question = "Quels sont les thèmes principaux de ce document ? Liste-les et explique chacun brièvement."
                    result, source_pages = ask_full_or_rag(client, question)
                    st.write(result)
                    st.caption(format_sources(source_pages))

    # TAB 4 : AUDIO
    with tabs[3]:
        max_page = len(st.session_state.pdf_pages)
        p_num = st.number_input("Numéro de page à lire", min_value=1, max_value=max_page, value=1)
        lang = st.selectbox("Langue", ["fr", "en", "es", "de"], index=0)

        if st.button("🔊 Générer l'audio", key="btn_audio"):
            page_text = st.session_state.pdf_pages.get(p_num, "")
            if page_text:
                with st.spinner("Génération audio..."):
                    try:
                        tts = gTTS(text=page_text, lang=lang)
                        audio_io = BytesIO()
                        tts.write_to_fp(audio_io)
                        audio_io.seek(0)
                        st.audio(audio_io, format="audio/mp3")
                        st.caption(f"📄 Page {p_num}")
                    except Exception as e:
                        st.error(f"Erreur audio : {e}")
            else:
                st.warning("Aucun texte trouvé sur cette page.")

    # TAB 5 : PRÉSENTATION
    with tabs[4]:
        n_slides = st.number_input("Nombre de slides", min_value=3, max_value=10, value=5)

        if st.button("🎯 Générer la présentation PPTX", key="btn_pptx"):
            with st.spinner("L'IA structure vos slides..."):
                try:
                    question = (
                        f"Crée une structure pour exactement {n_slides} slides basées sur ce document. "
                        f"Réponds UNIQUEMENT avec un JSON valide, sans texte avant ou après, sans balises markdown : "
                        f'{{ "slides": [ {{ "titre": "Titre de la slide", "points": ["Point 1", "Point 2", "Point 3"] }} ] }}'
                    )
                    raw, source_pages = ask_full_or_rag(client, question)
                    raw = raw.strip()
                    if raw.startswith("```"):
                        raw = re.sub(r"```(?:json)?", "", raw).strip("` \n")

                    start = raw.find('{')
                    end = raw.rfind('}') + 1

                    if start != -1 and end > 0:
                        data = json.loads(raw[start:end])
                        ppt_bytes = create_pptx(data)
                        st.download_button(
                            label="📥 Télécharger la présentation",
                            data=ppt_bytes,
                            file_name="presentation.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                        st.success(f"✅ {len(data.get('slides', []))} slides générées !")
                        st.caption(format_sources(source_pages))
                    else:
                        st.error("Format JSON invalide reçu.")
                        st.code(raw)
                except json.JSONDecodeError as e:
                    st.error(f"Erreur JSON : {e}")
                    st.code(raw)
                except Exception as e:
                    st.error(f"Erreur : {e}")

else:
    st.info("👈 Veuillez charger un fichier PDF dans la barre latérale pour commencer.")
